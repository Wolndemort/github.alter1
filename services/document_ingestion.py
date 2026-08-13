"""Safe text extraction for documents used by chat and durable agents."""
from __future__ import annotations

import io
import json
import re
from datetime import datetime
from dataclasses import dataclass
from pathlib import Path

from utils.agent_engine import start_agent
from utils.external_content import audit_external_content


MAX_DOCUMENT_BYTES = 25 * 1024 * 1024
MAX_DOCUMENT_CHARS = 120_000
SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".json", ".csv", ".rtf",
    ".pdf", ".docx", ".xlsx", ".pptx", ".odt",
}
EDITABLE_EXTENSIONS = set(SUPPORTED_EXTENSIONS)
OCR_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".tiff", ".bmp"}

DEFAULT_MEDIA_TYPES = {
    ".txt": "text/plain", ".md": "text/markdown", ".markdown": "text/markdown",
    ".csv": "text/csv", ".json": "application/json", ".rtf": "application/rtf",
    ".pdf": "application/pdf", ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".odt": "application/vnd.oasis.opendocument.text",
}


@dataclass(frozen=True)
class Document:
    filename: str
    media_type: str
    text: str
    pages: int | None = None
    ocr_used: bool = False

    @property
    def chars(self) -> int:
        return len(self.text)


@dataclass(frozen=True)
class EditedDocument:
    """A bounded, exportable document result produced by an explicit edit."""
    filename: str
    media_type: str
    data: bytes


def _clean(text: str) -> str:
    value = text.replace("\x00", "")
    value = re.sub(r"\n{3,}", "\n\n", value)
    value = re.sub(r"[ \t]+", " ", value)
    return value.strip()[:MAX_DOCUMENT_CHARS]


def _replace_document_text(text: str, old: str, new: str) -> tuple[str, int]:
    """Replace a phrase ignoring case and arbitrary whitespace/newlines."""
    parts = [re.escape(part) for part in old.split() if part]
    if not parts:
        return text, 0
    pattern = re.compile(r"\s+".join(parts), flags=re.IGNORECASE)
    return pattern.subn(lambda _: new, text)


def _extension(filename: str) -> str:
    return Path(filename or "document").suffix.casefold()


def _media_type(extension: str, supplied: str) -> str:
    return supplied or DEFAULT_MEDIA_TYPES.get(extension, "application/octet-stream")


def _extract_rtf(data: bytes) -> str:
    """Decode the readable text from common RTF without executing controls."""
    value = data.decode("latin1", errors="replace")
    result: list[str] = []
    index = 0
    skip_destination = 0
    while index < len(value):
        char = value[index]
        if char == "{":
            index += 1
            continue
        if char == "}":
            index += 1
            if skip_destination:
                skip_destination -= 1
            continue
        if char != "\\":
            if not skip_destination:
                result.append(char)
            index += 1
            continue
        index += 1
        if index >= len(value):
            break
        if value[index] in "\\{}":
            if not skip_destination:
                result.append(value[index])
            index += 1
            continue
        if value[index] == "'" and index + 2 < len(value):
            try:
                decoded = bytes.fromhex(value[index + 1:index + 3]).decode("cp1251")
            except (ValueError, UnicodeDecodeError):
                decoded = ""
            if not skip_destination:
                result.append(decoded)
            index += 3
            continue
        match = re.match(r"([a-zA-Z]+)(-?\d+)? ?", value[index:])
        if not match:
            index += 1
            continue
        word, number = match.groups()
        index += match.end()
        if word in {"par", "line", "tab"} and not skip_destination:
            result.append("\n" if word != "tab" else "\t")
        elif word == "u" and number and not skip_destination:
            codepoint = int(number)
            if codepoint < 0:
                codepoint += 65536
            result.append(chr(codepoint))
        elif word in {"fonttbl", "colortbl", "stylesheet", "info", "pict", "object"}:
            skip_destination += 1
    return "".join(result)


def _encode_rtf(text: str) -> bytes:
    escaped = []
    for char in text:
        if char == "\\":
            escaped.append(r"\\")
        elif char in "{}":
            escaped.append("\\" + char)
        elif char == "\n":
            escaped.append(r"\par ")
        elif ord(char) < 128:
            escaped.append(char)
        else:
            codepoint = ord(char)
            if codepoint > 32767:
                codepoint -= 65536
            escaped.append(f"\\u{codepoint}?" )
    return (r"{\rtf1\ansi\deff0 " + "".join(escaped) + "}").encode("ascii")


def extract_document(filename: str, data: bytes, media_type: str = "") -> Document:
    """Extract bounded text, raising a user-safe ValueError for unsupported input."""
    if not isinstance(data, (bytes, bytearray)) or not data:
        raise ValueError("document is empty")
    if len(data) > MAX_DOCUMENT_BYTES:
        raise ValueError("document is too large")
    filename = Path(filename or "document").name[:180]
    extension = _extension(filename)
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError("unsupported document type")
    if extension in {".txt", ".md", ".markdown", ".csv"}:
        text = bytes(data).decode("utf-8-sig", errors="replace")
        return Document(filename, _media_type(extension, media_type), _clean(text))
    if extension == ".rtf":
        return Document(filename, _media_type(extension, media_type), _clean(_extract_rtf(bytes(data))))
    if extension == ".json":
        try:
            value = json.loads(bytes(data).decode("utf-8-sig"))
            text = json.dumps(value, ensure_ascii=False, indent=2)
        except (UnicodeDecodeError, json.JSONDecodeError) as exc:
            raise ValueError("invalid JSON document") from exc
        return Document(filename, _media_type(extension, media_type), _clean(text))
    if extension == ".pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(bytes(data)))
            text = "\n\n".join(page.extract_text() or "" for page in reader.pages)
            return Document(filename, _media_type(extension, media_type), _clean(text), pages=len(reader.pages))
        except ImportError as exc:
            raise ValueError("PDF support is not installed") from exc
        except Exception as exc:
            raise ValueError("could not read PDF document") from exc
    if extension == ".xlsx":
        try:
            from openpyxl import load_workbook
            workbook = load_workbook(io.BytesIO(bytes(data)), read_only=True, data_only=False)
            parts = [f"[{sheet.title}]" for sheet in workbook.worksheets]
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    values = [str(value) for value in row if value is not None]
                    if values:
                        parts.append(" | ".join(values))
            workbook.close()
            return Document(filename, _media_type(extension, media_type), _clean("\n".join(parts)))
        except ImportError as exc:
            raise ValueError("XLSX support is not installed") from exc
        except Exception as exc:
            raise ValueError("could not read XLSX document") from exc
    if extension == ".pptx":
        try:
            from pptx import Presentation
            presentation = Presentation(io.BytesIO(bytes(data)))
            parts = []
            for number, slide in enumerate(presentation.slides, start=1):
                parts.append(f"[Slide {number}]")
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        parts.append(shape.text)
                    if getattr(shape, "has_table", False):
                        parts.extend(" | ".join(cell.text for cell in row.cells) for row in shape.table.rows)
            return Document(filename, _media_type(extension, media_type), _clean("\n".join(parts)), pages=len(presentation.slides))
        except ImportError as exc:
            raise ValueError("PPTX support is not installed") from exc
        except Exception as exc:
            raise ValueError("could not read PPTX document") from exc
    if extension == ".odt":
        try:
            from odf import teletype
            from odf.text import H, P
            from odf.opendocument import load
            source = load(io.BytesIO(bytes(data)))
            parts = [teletype.extractText(node) for node in source.getElementsByType(P)]
            parts.extend(teletype.extractText(node) for node in source.getElementsByType(H))
            return Document(filename, _media_type(extension, media_type), _clean("\n".join(parts)))
        except ImportError as exc:
            raise ValueError("ODT support is not installed") from exc
        except Exception as exc:
            raise ValueError("could not read ODT document") from exc
    try:
        from docx import Document as DocxDocument
        document = DocxDocument(io.BytesIO(bytes(data)))
        parts = [paragraph.text for paragraph in document.paragraphs]
        for table in document.tables:
            parts.extend(" | ".join(cell.text for cell in row.cells) for row in table.rows)
        return Document(filename, _media_type(extension, media_type), _clean("\n".join(parts)))
    except ImportError as exc:
        raise ValueError("DOCX support is not installed") from exc
    except Exception as exc:
        raise ValueError("could not read DOCX document") from exc


def document_chunks(document: Document, chunk_chars: int = 6000) -> list[str]:
    size = max(500, min(int(chunk_chars), 12000))
    return [document.text[index:index + size] for index in range(0, len(document.text), size)] or [""]


def document_profile(document: Document) -> dict:
    """Return bounded, deterministic signals for agents and audit UIs."""
    lines = [line.strip() for line in document.text.splitlines() if line.strip()]
    tables = [line for line in lines if line.count("|") >= 2 or "\t" in line]
    dates = sorted(set(re.findall(r"\b(?:\d{1,2}[./]){2}\d{2,4}\b|\b\d{4}-\d{2}-\d{2}\b", document.text)))
    amounts = sorted(set(re.findall(r"(?:\d[\d ]{0,12} ?(?:₽|руб\.?|\$|€|USD|EUR))", document.text, re.IGNORECASE)))
    return {
        "filename": document.filename,
        "media_type": document.media_type,
        "pages": document.pages,
        "chars": document.chars,
        "lines": len(lines),
        "tables": tables[:100],
        "dates": dates[:100],
        "amounts": amounts[:100],
        "ocr_used": document.ocr_used,
        "needs_ocr": document.pages is not None and not document.text,
    }


def ocr_image_text(filename: str, data: bytes, *, language: str = "rus+eng") -> Document:
    """Run bounded local OCR when Pillow/Tesseract are installed.

    OCR is deliberately optional: hosted vision remains the fallback when a
    deployment does not have the native Tesseract binary.
    """
    if not data:
        raise ValueError("image is empty")
    if len(data) > MAX_DOCUMENT_BYTES:
        raise ValueError("image is too large")
    if _extension(filename) not in OCR_EXTENSIONS:
        raise ValueError("unsupported OCR image type")
    try:
        from PIL import Image
        import pytesseract
        image = Image.open(io.BytesIO(bytes(data)))
        text = pytesseract.image_to_string(image, lang=language)
    except ImportError as exc:
        raise ValueError("local OCR is unavailable; use vision analysis") from exc
    except Exception as exc:
        raise ValueError("could not OCR image") from exc
    return Document(Path(filename).name[:180], "text/plain", _clean(text), ocr_used=True)


def edit_document(filename: str, data: bytes, instruction: str, media_type: str = "") -> EditedDocument:
    """Apply only deterministic replacements supplied as ``old => new`` lines.

    AI-generated edits must be converted to this explicit format by the caller;
    this keeps exports auditable and prevents accidental destructive rewrites.
    """
    extension = _extension(filename)
    if extension not in EDITABLE_EXTENSIONS:
        raise ValueError("this document format requires layout-aware editing")
    replacements = []
    for line in (instruction or "").splitlines():
        if "=>" in line:
            old, new = line.split("=>", 1)
            if old.strip():
                replacements.append((old.strip(), new.strip()))
    if not replacements:
        raise ValueError("provide explicit replacements in the form: old => new")
    if extension == ".pdf":
        return edit_pdf_document(filename, data, replacements, media_type)
    document = extract_document(filename, data, media_type)
    text = document.text
    for old, new in replacements:
        text, count = _replace_document_text(text, old, new)
        if count == 0:
            raise ValueError("не нашёл в документе текст для замены: " + repr(old[:80]))
    if extension == ".json":
        try:
            output = json.dumps(json.loads(text), ensure_ascii=False, indent=2).encode("utf-8")
        except json.JSONDecodeError as exc:
            raise ValueError("edit would produce invalid JSON") from exc
    elif extension == ".docx":
        try:
            from docx import Document as DocxDocument
            source = DocxDocument(io.BytesIO(bytes(data)))
            changed_by_replacement = [0 for _ in replacements]

            def replace_paragraph(paragraph) -> None:
                value = paragraph.text
                for index, (old, new) in enumerate(replacements):
                    value, count = _replace_document_text(value, old, new)
                    changed_by_replacement[index] += count
                if value != paragraph.text:
                    if paragraph.runs:
                        paragraph.runs[0].text = value
                        for run in paragraph.runs[1:]:
                            run.text = ""
                    else:
                        paragraph.add_run(value)

            for paragraph in source.paragraphs:
                replace_paragraph(paragraph)
            for table in source.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            replace_paragraph(paragraph)
            for (old, _), count in zip(replacements, changed_by_replacement):
                if count == 0:
                    raise ValueError("не нашёл в документе текст для замены: " + repr(old[:80]))
            output_buffer = io.BytesIO()
            source.save(output_buffer)
            output = output_buffer.getvalue()
        except ImportError as exc:
            raise ValueError("DOCX support is not installed") from exc
    elif extension == ".xlsx":
        try:
            from openpyxl import load_workbook
            source = load_workbook(io.BytesIO(bytes(data)))
            changed_by_replacement = [0 for _ in replacements]
            for sheet in source.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if not isinstance(cell.value, str):
                            continue
                        value = cell.value
                        for index, (old, new) in enumerate(replacements):
                            value, count = _replace_document_text(value, old, new)
                            changed_by_replacement[index] += count
                        cell.value = value
            source_buffer = io.BytesIO()
            source.save(source_buffer)
            source.close()
            for (old, _), count in zip(replacements, changed_by_replacement):
                if count == 0:
                    raise ValueError("не нашёл в документе текст для замены: " + repr(old[:80]))
            output = source_buffer.getvalue()
        except ImportError as exc:
            raise ValueError("XLSX support is not installed") from exc
    elif extension == ".pptx":
        try:
            from pptx import Presentation
            source = Presentation(io.BytesIO(bytes(data)))
            changed_by_replacement = [0 for _ in replacements]

            def replace_shape(shape) -> None:
                if getattr(shape, "has_text_frame", False):
                    value = shape.text
                    for index, (old, new) in enumerate(replacements):
                        value, count = _replace_document_text(value, old, new)
                        changed_by_replacement[index] += count
                    shape.text = value
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            value = cell.text
                            for index, (old, new) in enumerate(replacements):
                                value, count = _replace_document_text(value, old, new)
                                changed_by_replacement[index] += count
                            cell.text = value
                if hasattr(shape, "shapes"):
                    for child in shape.shapes:
                        replace_shape(child)

            for slide in source.slides:
                for shape in slide.shapes:
                    replace_shape(shape)
            for (old, _), count in zip(replacements, changed_by_replacement):
                if count == 0:
                    raise ValueError("не нашёл в документе текст для замены: " + repr(old[:80]))
            output_buffer = io.BytesIO()
            source.save(output_buffer)
            output = output_buffer.getvalue()
        except ImportError as exc:
            raise ValueError("PPTX support is not installed") from exc
    elif extension == ".odt":
        try:
            from odf import teletype
            from odf.opendocument import load
            from odf.text import H, P
            source = load(io.BytesIO(bytes(data)))
            changed_by_replacement = [0 for _ in replacements]
            for node in source.getElementsByType(P) + source.getElementsByType(H):
                value = teletype.extractText(node)
                for index, (old, new) in enumerate(replacements):
                    value, count = _replace_document_text(value, old, new)
                    changed_by_replacement[index] += count
                if value != teletype.extractText(node):
                    # ODF text nodes are not removable through Element's
                    # removeChild cache path. Replacing the bounded paragraph
                    # contents also avoids leaving stale formatted runs.
                    node.childNodes[:] = []
                    node.addText(value)
            for (old, _), count in zip(replacements, changed_by_replacement):
                if count == 0:
                    raise ValueError("не нашёл в документе текст для замены: " + repr(old[:80]))
            output_buffer = io.BytesIO()
            source.save(output_buffer)
            output = output_buffer.getvalue()
        except ImportError as exc:
            raise ValueError("ODT support is not installed") from exc
    elif extension == ".rtf":
        output = _encode_rtf(text)
    else:
        output = text.encode("utf-8")
    return EditedDocument(document.filename, document.media_type, output)


def edit_pdf_document(filename: str, data: bytes, replacements: list[tuple[str, str]], media_type: str = "") -> EditedDocument:
    """Apply coordinate-aware redactions to text-layer PDFs.

    ``page.search_for`` is fast but provider/PDF text often differs in case
    or is split across lines.  The word-based fallback below makes PDF edits
    use the same tolerant matching contract as TXT/DOCX while preserving the
    original page layout.  A scanned PDF has no words and is rejected clearly
    instead of returning a falsely successful unchanged file.
    """
    try:
        import fitz
    except ImportError as exc:
        raise ValueError("PDF layout editing is not installed") from exc
    def search_rects(page, phrase: str) -> list[object]:
        wanted = re.findall(r"\w+", phrase.casefold(), flags=re.UNICODE)
        if not wanted:
            return []
        words = page.get_text("words") or []
        actual = [str(item[4]) for item in words if len(item) >= 5]
        actual_tokens = [re.findall(r"\w+", item.casefold(), flags=re.UNICODE) for item in actual]
        flat: list[tuple[int, str]] = []
        for index, tokens in enumerate(actual_tokens):
            flat.extend((index, token) for token in tokens)
        matches: list[object] = []
        for start in range(0, len(flat) - len(wanted) + 1):
            if [token for _, token in flat[start:start + len(wanted)]] != wanted:
                continue
            indexes = {index for index, _ in flat[start:start + len(wanted)]}
            rect = None
            for index in indexes:
                item = words[index]
                current = fitz.Rect(item[:4])
                rect = current if rect is None else rect | current
            if rect is not None:
                matches.append(rect)
        return matches

    try:
        source = fitz.open(stream=bytes(data), filetype="pdf")
        changed = 0
        for page in source:
            for old, new in replacements:
                for rect in search_rects(page, old):
                    page.add_redact_annot(rect, text=new, fontname="helv", fontsize=max(6, min(18, rect.height * 0.8)), align=0)
                    changed += 1
            page.apply_redactions()
        if not changed:
            source.close()
            raise ValueError("PDF text was not found; scanned PDFs require OCR before editing")
        output = source.tobytes(garbage=4, deflate=True)
        source.close()
        return EditedDocument(Path(filename or "document.pdf").name[:180], media_type or "application/pdf", output)
    except ValueError:
        raise
    except Exception as exc:
        raise ValueError("could not edit PDF layout") from exc


def start_document_agent(settings: dict | None, document: Document, goal: str, *, horizon_minutes: int = 60) -> dict:
    """Create a durable document agent without storing unbounded raw input."""
    bounded_text = document.text[:30000]
    tasks = [
        {"id": "document_scope", "title": "Определить цель и структуру документа", "priority": 1},
        {"id": "document_facts", "title": "Извлечь ключевые факты, даты, числа и ограничения", "depends_on": ["document_scope"]},
        {"id": "document_actions", "title": "Сформировать практические выводы и задачи", "depends_on": ["document_facts"]},
        {"id": "document_verify", "title": "Проверить выводы и подготовить результат", "depends_on": ["document_actions"]},
    ]
    return start_agent(
        settings, goal or f"Работа с документом {document.filename}", horizon_minutes=horizon_minutes,
        tasks=tasks, constraints={"document_context": bounded_text, "document_filename": document.filename, "document_profile": document_profile(document), "external_content_audit": audit_external_content(bounded_text)},
    )
